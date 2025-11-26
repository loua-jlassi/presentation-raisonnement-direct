#!/usr/bin/env python3
"""
Script pour créer directement un fichier PowerPoint (.pptx) à partir de la présentation
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Installation de python-pptx...")
    import subprocess
    import sys
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "--quiet"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

# Couleurs
PRIMARY = RGBColor(102, 126, 234)  # #667eea
SECONDARY = RGBColor(118, 75, 162)  # #764ba2
SUCCESS = RGBColor(39, 174, 96)  # #27ae60
WARNING = RGBColor(255, 152, 0)  # #ff9800
INFO = RGBColor(33, 150, 243)  # #2196f3
DANGER = RGBColor(231, 76, 60)  # #e74c3c

def add_title_slide(prs, title, subtitle, author, course, year):
    """Ajoute une slide de titre"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout vide
    
    # Fond avec gradient (simulé avec un rectangle)
    from pptx.shapes.autoshape import Shape
    left = top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height
    bg = slide.shapes.add_shape(1, left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(2))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Sous-titre
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(10), Inches(1))
    tf = subtitle_box.text_frame
    tf.text = subtitle
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Informations
    info_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(10), Inches(2))
    tf = info_box.text_frame
    tf.text = f"Présenté par : {author}\nCours : {course}\nAnnée académique : {year}"
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_items):
    """Ajoute une slide de contenu"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Titre et contenu
    
    # Titre
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    
    # Contenu - utiliser le placeholder de contenu
    if len(slide.placeholders) > 1:
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i > 0:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(18)
            p.level = 0
    else:
        # Si pas de placeholder, créer une textbox
        left = Inches(0.5)
        top = Inches(2)
        width = Inches(9)
        height = Inches(5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i > 0:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(18)
            p.level = 0

def add_bullet_slide(prs, title, bullets):
    """Ajoute une slide avec puces"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Titre et contenu
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    
    content_shape = slide.placeholders[1]
    tf = content_shape.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = bullet
        p.font.size = Pt(20)
        p.level = 0

def create_presentation():
    """Crée la présentation PowerPoint complète"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 0: Titre
    add_title_slide(
        prs,
        "RAISONNEMENT DIRECT",
        "Méthodes et applications du raisonnement logique",
        "Loua El Jelassi",
        "Représentation de connaissances et raisonnements",
        "2025 / 2026"
    )
    
    # Slide 1: Plan
    add_bullet_slide(prs, "Plan de la Présentation", [
        "1. Introduction - Historique",
        "2. Qu'est-ce que le Raisonnement Direct ?",
        "3. Architecture du Raisonnement Direct",
        "4. Méthodes d'Inférence",
        "5. Modus Ponens",
        "6. Chaînage Avant",
        "7. Chaînage Arrière",
        "8. Exemple d'Application",
        "9. Avantages & Limites",
        "10. Conclusion"
    ])
    
    # Slide 2: Introduction
    add_content_slide(prs, "Introduction", [
        "Père de la logique formelle : Aristote (384–322 av. J.-C.)",
        "Il introduit le syllogisme dans Organon, base du raisonnement déductif moderne."
    ])
    
    # Slide 3: Évolution Historique
    add_bullet_slide(prs, "L'Évolution Historique", [
        "XXᵉ siècle : Frege, Tarski et Turing formalisent la logique",
        "Années 1950–1970 : Naissance des systèmes experts et des moteurs d'inférence",
        "Aujourd'hui : Utilisé dans l'IA, la vérification formelle et les bases de données déductives"
    ])
    
    # Slide 4: Définition
    add_content_slide(prs, "Qu'est-ce que le raisonnement direct ?", [
        "Le raisonnement direct est une méthode de déduction logique qui permet d'inférer de nouvelles connaissances à partir de faits et de règles déjà connus.",
        "Il repose sur l'application systématique et déterministe de règles d'inférence afin de produire des conclusions logiquement valides."
    ])
    
    # Slide 5: Principe de base
    add_bullet_slide(prs, "Principe de base", [
        "1. Prend des faits initiaux (prémisses)",
        "2. Applique des règles logiques (règles d'inférence)",
        "3. Génère de nouvelles conclusions (faits dérivés)"
    ])
    
    # Slide 6: Schéma du processus
    add_content_slide(prs, "Schéma du processus", [
        "FAITS INITIAUX (P, Q, R...) → RÈGLES (P→Q, Q→R...) → CONCLUSIONS (S, T, U...)"
    ])
    
    # Slide 7: Caractéristiques
    add_bullet_slide(prs, "Caractéristiques du raisonnement direct", [
        "Déterministe : Produit toujours le même résultat pour les mêmes prémisses",
        "Monotone : Les conclusions restent valides même après l'ajout de nouvelles informations",
        "Complet : Toutes les conclusions logiquement possibles peuvent être déduites",
        "Sûr : Ne produit que des conclusions logiquement valides"
    ])
    
    # Slide 8: Types
    add_bullet_slide(prs, "Types de raisonnement direct", [
        "Chaînage avant : Faits → Conclusions",
        "Chaînage arrière : But → Prémisses"
    ])
    
    # Slide 9: Architecture - Composants
    add_bullet_slide(prs, "Architecture - Composants principaux", [
        "1. Base de faits : Contient les connaissances connues",
        "2. Base de règles : Regroupe les règles d'inférence",
        "3. Moteur d'inférence : Cœur du système",
        "4. Mémoire de travail : Espace temporaire pour les conclusions intermédiaires",
        "5. Interface utilisateur : Permet l'interaction avec le système"
    ])
    
    # Slide 10: Architecture - Schéma
    add_content_slide(prs, "Schéma simplifié de l'architecture", [
        "Base de Règles ↑",
        "Interface/Utilisateur → Moteur d'Inférence → Base de Faits",
        "↓ Mémoire de Travail"
    ])
    
    # Slide 11: Fonctionnement global
    add_bullet_slide(prs, "Fonctionnement global", [
        "1. L'utilisateur saisit les faits initiaux",
        "2. Le moteur d'inférence identifie les règles applicables",
        "3. Les conclusions déduites sont ajoutées à la base de faits",
        "4. Le processus se répète jusqu'à ce qu'aucune règle ne soit plus applicable",
        "5. Le module d'explication présente les résultats finaux"
    ])
    
    # Slide 12: Méthodes d'inférence
    add_bullet_slide(prs, "Les Méthodes d'Inférence les Plus Utilisées", [
        "Modus Ponens : Règle d'inférence de base",
        "Chaînage avant : Part des faits connus",
        "Chaînage arrière : Part d'un but à vérifier"
    ])
    
    # Slide 13: Modus Ponens
    add_content_slide(prs, "Modus Ponens (Règle Fondamentale)", [
        "Principe : Règle d'inférence fondamentale sur laquelle reposent toutes les autres méthodes",
        "Forme générale :",
        "  P → Q",
        "  P",
        "  ∴ Q",
        "En mots simples : Si P implique Q et que P est vrai, alors Q est nécessairement vrai."
    ])
    
    # Slide 14: Modus Ponens - Avantages
    add_bullet_slide(prs, "Modus Ponens - Avantages", [
        "Simple à comprendre et à appliquer",
        "Logique et universelle",
        "Sert de base aux chaînages avant et arrière",
        "Utilisée dans tous les moteurs d'inférence"
    ])
    
    # Slide 15: Chaînage Avant
    add_content_slide(prs, "Chaînage Avant (Forward Chaining)", [
        "Principe : Méthode la plus utilisée dans les systèmes experts",
        "Fonctionnement :",
        "1. Le moteur utilise la base de faits initiale",
        "2. Il applique toutes les règles dont les conditions sont satisfaites",
        "3. Les conclusions sont ajoutées à la base de faits",
        "4. Le processus se répète jusqu'à épuisement"
    ])
    
    # Slide 16: Chaînage Avant - Exemple
    add_content_slide(prs, "Chaînage Avant - Exemple", [
        "Règle : S'il pleut → la route est mouillée",
        "Fait : Il pleut",
        "⇒ Conclusion : La route est mouillée",
        "",
        "Avantages :",
        "- Facile à implémenter",
        "- Très efficace pour générer de nouvelles connaissances"
    ])
    
    # Slide 17: Chaînage Arrière
    add_content_slide(prs, "Chaînage Arrière (Backward Chaining)", [
        "Principe : Part d'un objectif ou d'une conclusion à vérifier et remonte pour déterminer les faits nécessaires",
        "Fonctionnement :",
        "1. L'utilisateur définit un objectif",
        "2. Le moteur identifie les règles pouvant produire cette conclusion",
        "3. Il vérifie si les conditions sont satisfaites",
        "4. Si non, il cherche les faits nécessaires",
        "5. Le processus se répète jusqu'à validation ou échec"
    ])
    
    # Slide 18: Chaînage Arrière - Exemple
    add_content_slide(prs, "Chaînage Arrière - Exemple", [
        "Objectif : La route est mouillée",
        "Règle : S'il pleut → la route est mouillée",
        "Fait : Il pleut",
        "⇒ Conclusion : L'objectif « la route est mouillée » est validé",
        "",
        "Avantages :",
        "- Permet de se concentrer sur un objectif précis",
        "- Efficace pour les systèmes de diagnostic et de planification"
    ])
    
    # Slide 19: Exemple d'application
    add_content_slide(prs, "Exemple d'application - Système de diagnostic médical", [
        "Fonctionnement :",
        "1. Le patient présente des symptômes (fièvre, toux, fatigue)",
        "2. Le système applique les règles de diagnostic",
        "3. Il génère un diagnostic avec probabilités et recommandations"
    ])
    
    # Slide 20: Exemple - Schéma
    add_content_slide(prs, "Schéma du processus de diagnostic", [
        "Base de faits : Fièvre • Toux • Fatigue",
        "↓",
        "Moteur d'inférence : Application des règles",
        "  Règle 1: Fièvre + Toux → Grippe",
        "  Règle 2: Fièvre + Fatigue → Infection",
        "↓",
        "Diagnostic : Grippe (probabilité: 85%)",
        "Traitement: Repos, hydratation, médicaments"
    ])
    
    # Slide 21: Avantages
    add_bullet_slide(prs, "Principaux avantages", [
        "Déterminisme : Résultats prévisibles et reproductibles",
        "Sûreté : Seules des conclusions logiquement valides sont produites",
        "Efficacité : Rapide et simple à mettre en œuvre",
        "Transparence : Raisonnement traçable et compréhensible"
    ])
    
    # Slide 22: Limites
    add_bullet_slide(prs, "Principales limites", [
        "Rigidité : Ne gère pas l'incertitude ou les informations incomplètes",
        "Complexité computationnelle : Peut être coûteux pour de grandes bases",
        "Absence d'apprentissage : Ne s'améliore pas automatiquement",
        "Explosion combinatoire : Le nombre de règles peut croître rapidement"
    ])
    
    # Slide 23: Conclusion
    add_content_slide(prs, "Conclusion", [
        "Le raisonnement direct permet d'inférer de nouvelles connaissances de façon logique, prévisible et traçable.",
        "",
        "Il est largement utilisé dans les systèmes experts et l'intelligence artificielle, malgré ses limites.",
        "",
        "C'est un outil fondamental pour l'automatisation du raisonnement logique et il continue d'évoluer avec les technologies modernes."
    ])
    
    # Slide 24: Questions
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond
    left = top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height
    bg = slide.shapes.add_shape(1, left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(1.5))
    tf = title_box.text_frame
    tf.text = "Questions et discussion"
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Message
    msg_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(10), Inches(2))
    tf = msg_box.text_frame
    tf.text = "Merci pour votre attention !\n\nPrêt à répondre à vos questions sur le raisonnement direct et ses applications."
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(28)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER
    
    return prs

def main():
    print("=" * 60)
    print("CRÉATION DU FICHIER POWERPOINT (.pptx)")
    print("=" * 60)
    print()
    
    try:
        prs = create_presentation()
        output_file = "presentation_raisonnement_direct.pptx"
        prs.save(output_file)
        print(f"✓ Présentation créée avec succès : {output_file}")
        print(f"✓ Nombre de slides : {len(prs.slides)}")
        print()
        print("Vous pouvez maintenant ouvrir ce fichier dans PowerPoint !")
        print()
        
        # Ouvrir le fichier
        import os
        import subprocess
        file_path = os.path.abspath(output_file)
        print(f"Ouverture de {file_path}...")
        subprocess.run(["open", file_path], check=False)
        
    except Exception as e:
        print(f"❌ Erreur : {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    main()

